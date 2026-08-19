# -*- coding: utf-8 -*-
"""
文档生成工具集 - 支持 PPT、PDF、Word、Excel 生成
"""

import os
import json
from pathlib import Path
from datetime import datetime


def _import_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.dml.color import RGBColor
        return Presentation, Inches, Pt, Emu, PP_ALIGN, MSO_ANCHOR, RGBColor
    except ImportError:
        return None


def _import_docx():
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        return Document, Inches, Pt, RGBColor, WD_ALIGN_PARAGRAPH
    except ImportError:
        return None


def _import_openpyxl():
    try:
        from openpyxl import Workbook, load_workbook
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        from openpyxl.utils import get_column_letter
        return Workbook, load_workbook, Font, Alignment, PatternFill, Border, Side, get_column_letter
    except ImportError:
        return None


def _import_reportlab():
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.lib.units import inch
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        return canvas, letter, A4, inch, pdfmetrics, TTFont
    except ImportError:
        return None


def _ensure_dir(path):
    dirname = os.path.dirname(path)
    if dirname and not os.path.exists(dirname):
        os.makedirs(dirname, exist_ok=True)


def _get_extension(path):
    return os.path.splitext(path)[1].lower()


def _generate_ppt(args: dict) -> str:
    try:
        path = args.get('path', 'output.pptx')
        title = args.get('title', '演示文稿')
        slides = args.get('slides', [])
        author = args.get('author', 'Hermes')
        _ensure_dir(path)
        pptx_module = _import_pptx()
        if not pptx_module:
            return "❌ 请安装 python-pptx: pip install python-pptx"
        Presentation, Inches, Pt, Emu, PP_ALIGN, MSO_ANCHOR, RGBColor = pptx_module
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 30, 30)
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        for slide_data in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide_title = slide_data.get('title', '')
            content = slide_data.get('content', [])
            if isinstance(content, str):
                content = content.split('\n')
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1))
            tf = title_box.text_frame
            tf.text = slide_title
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11.333), Inches(5.5))
            tf = text_box.text_frame
            tf.word_wrap = True
            for item in content:
                if item.strip():
                    p = tf.add_paragraph()
                    p.text = item
                    p.font.size = Pt(24)
        prs.save(path)
        return f"✅ PPT 已生成: {path} ({len(slides) + 1} 页)"
    except Exception as e:
        return f"❌ 生成 PPT 失败: {e}"


def _generate_pdf(args: dict) -> str:
    try:
        path = args.get('path', 'output.pdf')
        content = args.get('content', '')
        title = args.get('title', '文档')
        author = args.get('author', 'Hermes')
        _ensure_dir(path)
        reportlab_module = _import_reportlab()
        if not reportlab_module:
            return "❌ 请安装 reportlab: pip install reportlab"
        canvas_mod, letter, A4, inch, pdfmetrics, TTFont = reportlab_module
        font_name = 'Helvetica'
        font_paths = ['C:/Windows/Fonts/simsun.ttc', 'C:/Windows/Fonts/simhei.ttf', 'C:/Windows/Fonts/msyh.ttc']
        for fp in font_paths:
            if os.path.exists(fp):
                try:
                    pdfmetrics.registerFont(TTFont('ChineseFont', fp))
                    font_name = 'ChineseFont'
                    break
                except:
                    continue
        c = canvas_mod.Canvas(path, pagesize=A4)
        width, height = A4
        c.setFont(font_name, 24)
        c.drawString(1 * inch, height - 1 * inch, title)
        c.setFont(font_name, 12)
        y = height - 2.2 * inch
        for line in content.split('\n'):
            if y < 1 * inch:
                c.showPage()
                c.setFont(font_name, 12)
                y = height - 1 * inch
            if line.strip():
                c.drawString(1 * inch, y, line[:80])
                y -= 0.3 * inch
        c.save()
        return f"✅ PDF 已生成: {path}"
    except Exception as e:
        return f"❌ 生成 PDF 失败: {e}"


def _generate_word(args: dict) -> str:
    try:
        path = args.get('path', 'output.docx')
        content = args.get('content', '')
        title = args.get('title', '文档')
        author = args.get('author', 'Hermes')
        _ensure_dir(path)
        docx_module = _import_docx()
        if not docx_module:
            return "❌ 请安装 python-docx: pip install python-docx"
        Document, Inches, Pt, RGBColor, WD_ALIGN_PARAGRAPH = docx_module
        doc = Document()
        heading = doc.add_heading(title, 0)
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for line in content.split('\n'):
            if line.strip():
                if line.startswith('# '):
                    doc.add_heading(line[2:], 1)
                elif line.startswith('## '):
                    doc.add_heading(line[3:], 2)
                else:
                    doc.add_paragraph(line)
            else:
                doc.add_paragraph()
        doc.save(path)
        return f"✅ Word 文档已生成: {path}"
    except Exception as e:
        return f"❌ 生成 Word 失败: {e}"


def _generate_excel(args: dict) -> str:
    try:
        path = args.get('path', 'output.xlsx')
        sheets = args.get('sheets', [])
        if not sheets:
            sheets = [{'name': 'Sheet1', 'data': args.get('data', [])}]
        _ensure_dir(path)
        openpyxl_module = _import_openpyxl()
        if not openpyxl_module:
            return "❌ 请安装 openpyxl: pip install openpyxl"
        Workbook, load_workbook, Font, Alignment, PatternFill, Border, Side, get_column_letter = openpyxl_module
        wb = Workbook()
        for idx, sheet_data in enumerate(sheets):
            if idx == 0:
                ws = wb.active
                ws.title = sheet_data.get('name', f'Sheet{idx+1}')[:31]
            else:
                ws = wb.create_sheet(sheet_data.get('name', f'Sheet{idx+1}')[:31])
            data_rows = sheet_data.get('data', [])
            for row_idx, row in enumerate(data_rows, 1):
                for col_idx, value in enumerate(row, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    if row_idx == 1:
                        cell.font = Font(bold=True)
        wb.save(path)
        return f"✅ Excel 已生成: {path}"
    except Exception as e:
        return f"❌ 生成 Excel 失败: {e}"


def register_tools():
    """注册工具到 Hermes"""
    import sys
    sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    import tools

    tools.register(name="generate_ppt", description="生成PowerPoint演示文稿。参数: path, title, slides, author", parameters={"type": "object", "properties": {"path": {"type": "string"}, "title": {"type": "string"}, "slides": {"type": "array"}, "author": {"type": "string"}}}, func=_generate_ppt)
    tools.register(name="generate_pdf", description="生成PDF文件。参数: path, content, title, author", parameters={"type": "object", "properties": {"path": {"type": "string"}, "content": {"type": "string"}, "title": {"type": "string"}, "author": {"type": "string"}}}, func=_generate_pdf)
    tools.register(name="generate_word", description="生成Word文档。参数: path, content, title, author", parameters={"type": "object", "properties": {"path": {"type": "string"}, "content": {"type": "string"}, "title": {"type": "string"}, "author": {"type": "string"}}}, func=_generate_word)
    tools.register(name="generate_excel", description="生成Excel表格。参数: path, sheets, data", parameters={"type": "object", "properties": {"path": {"type": "string"}, "sheets": {"type": "array"}, "data": {"type": "array"}}}, func=_generate_excel)
    return 4


def unregister_tools():
    """卸载工具"""
    import sys
    sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    import tools
    for name in ["generate_ppt", "generate_pdf", "generate_word", "generate_excel"]:
        tools.TOOLS.pop(name, None)
